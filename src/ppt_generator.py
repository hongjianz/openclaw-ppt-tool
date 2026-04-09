"""
PPT生成核心模块 - 使用python-pptx生成演示文稿
支持自动分页、排版和背景模板应用
"""

import os
import re
import urllib.request
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .template_config import TemplateConfig
from .content_parser import SlideContent, PresentationContent, TableData, CodeBlock, ChartData
from .smart_layout import create_layout_plan, ContentBlock


def hex_to_rgb(hex_color: str) -> RGBColor:
    """将十六进制颜色转换为RGBColor对象"""
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def safe_color(color_str: str, default: str = "#333333") -> RGBColor:
    """
    安全地解析颜色字符串
    如果不支持则使用默认值
    """
    if not color_str:
        return hex_to_rgb(default)

    # 处理rgba/rgb格式 - 提取主要颜色或使用默认值
    if color_str.startswith('rgba') or color_str.startswith('rgb'):
        # 对于QuintenStyle,白色文本使用白色
        if '255,255,255' in color_str or '255, 255, 255' in color_str:
            return RGBColor(255, 255, 255)
        return hex_to_rgb(default)

    # 标准hex颜色
    if color_str.startswith('#') and len(color_str) == 7:
        try:
            return hex_to_rgb(color_str)
        except (ValueError, IndexError):
            return hex_to_rgb(default)

    return hex_to_rgb(default)


def apply_bold_formatting(paragraph, text: str):
    """
    应用Markdown加粗格式到段落

    Args:
        paragraph: pptx段落对象
        text: 可能包含**text**语法的文本
    """
    # 检测是否包含加粗语法
    if '**' not in text and '__' not in text:
        # 没有加粗语法，直接设置文本
        paragraph.text = text
        return

    # 清除现有内容
    for run in paragraph.runs:
        run.text = ""

    # 使用正则表达式分割文本
    import re
    parts = re.split(r'(\*\*.*?\*\*|__.*?__)', text)

    for part in parts:
        if not part:
            continue

        # 检查是否为加粗部分
        is_bold = False
        clean_text = part

        if part.startswith('**') and part.endswith('**'):
            is_bold = True
            clean_text = part[2:-2]
        elif part.startswith('__') and part.endswith('__'):
            is_bold = True
            clean_text = part[2:-2]

        # 添加run
        run = paragraph.add_run()
        run.text = clean_text
        run.font.bold = is_bold


def calculate_text_height(text: str, font_size_pt: float, line_spacing: float = 1.2, max_width_inches: float = 10.0) -> float:
    """
    计算文本实际占用高度（英寸）

    Args:
        text: 文本内容
        font_size_pt: 字体大小（Pt）
        line_spacing: 行间距倍数（默认1.2）
        max_width_inches: 最大宽度（英寸）

    Returns:
        文本占用的实际高度（英寸）
    """
    if not text:
        return 0.0

    # 更准确的文本高度估算
    # 中文字体：约10字/英寸（18pt字体）
    # 英文字体：约15-20字符/英寸
    # 根据字体大小调整
    base_chars_per_inch = 10  # 中文基准
    scale_factor = 18.0 / font_size_pt  # 字体越小，每行字符越多
    chars_per_line = int(max_width_inches * base_chars_per_inch * scale_factor)
    chars_per_line = max(chars_per_line, 1)

    # 按换行符分割，计算实际行数
    lines = text.split('\n')
    total_lines = 0
    for line in lines:
        if not line.strip():
            total_lines += 1  # 空行也算一行
        else:
            # 计算该行需要的行数
            line_chars = len(line)
            lines_for_this = max(1, (line_chars + chars_per_line - 1) // chars_per_line)
            total_lines += lines_for_this

    # 计算总高度：行数 * (字体大小 / 72 * 行间距)
    line_height_inches = (font_size_pt / 72.0) * line_spacing
    total_height = total_lines * line_height_inches

    # 添加额外边距（段落间距）
    paragraph_count = len([l for l in lines if l.strip()])
    total_height += paragraph_count * 0.1  # 每段额外0.1英寸

    return total_height


class PPTGenerator:
    """PPT生成器"""

    def __init__(self, config: TemplateConfig):
        self.config = config
        self.prs = Presentation()
        # 设置页面尺寸
        self.prs.slide_width = Inches(config.slide_width)
        self.prs.slide_height = Inches(config.slide_height)

    def set_background(self, slide):
        """为幻灯片设置背景"""
        # 优先使用背景图片
        if self.config.background_image and os.path.exists(self.config.background_image):
            left = top = Inches(0)
            width = self.prs.slide_width
            height = self.prs.slide_height
            slide.shapes.add_picture(
                self.config.background_image, left, top,
                width=width, height=height
            )
        # QuintenStyle特殊处理:自动生成背景图片
        elif self.config.use_gradient_background and self.config.style_name == "QuintenStyle":
            # 尝试生成背景图片
            bg_path = self._generate_quinten_background()
            if bg_path and os.path.exists(bg_path):
                left = top = Inches(0)
                width = self.prs.slide_width
                height = self.prs.slide_height
                slide.shapes.add_picture(
                    bg_path, left, top,
                    width=width, height=height
                )
            else:
                # 降级为纯色
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = hex_to_rgb(self.config.background_color)
        else:
            # 使用纯色背景
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(self.config.background_color)

    def _generate_quinten_background(self):
        """生成或获取QuintenStyle背景图片"""
        # 检查templates目录
        templates_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'templates')
        bg_path = os.path.join(templates_dir, 'quinten_background.png')

        # 如果已存在,直接返回
        if os.path.exists(bg_path):
            return bg_path

        # 尝试生成
        try:
            from generate_quinten_background import generate_quinten_background
            os.makedirs(templates_dir, exist_ok=True)
            generate_quinten_background(1920, 1080, bg_path)
            return bg_path
        except Exception as e:
            print(f"警告: 无法生成QuintenStyle背景图片: {e}")
            return None

    def _add_table(self, slide, table_data: TableData, left, top, width) -> float:
        """
        添加表格到幻灯片

        Args:
            slide: 幻灯片对象
            table_data: 表格数据
            left: 左边距
            top: 顶部位置
            width: 表格宽度

        Returns:
            表格实际占用高度（英寸）
        """
        rows = len(table_data.rows) + 1  # +1 for header
        cols = len(table_data.headers)

        if rows == 0 or cols == 0:
            return 0.0

        # 计算行高（统一使用浮点数）
        row_height_inches = 0.5  # 每行0.5英寸
        height_inches = row_height_inches * rows

        # 添加表格（API调用时转为Inches对象）
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(height_inches))
        table = table_shape.table

        # 设置表头
        for i, header in enumerate(table_data.headers):
            cell = table.cell(0, i)
            cell.text = header

            # 设置表头样式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(self.config.body_size - 2)
                paragraph.font.bold = True
                paragraph.font.color.rgb = safe_color("#1A4D4D", "#FFFFFF")
                paragraph.alignment = PP_ALIGN.CENTER

            # 设置表头背景（使用纯色，因为transparency可能不支持）
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = safe_color("#B0C5BE", "#B0C5BE")
            # 注意：fill.transparency 在某些python-pptx版本中可能不可用
            # 如需透明度，建议使用浅色背景代替

        # 设置数据行
        for row_idx, row_data in enumerate(table_data.rows, 1):
            for col_idx, cell_data in enumerate(row_data):
                if col_idx < cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = cell_data

                    # 设置单元格样式
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(self.config.body_size - 2)
                        paragraph.font.color.rgb = safe_color("#1A4D4D", "#333333")

                        # 设置对齐方式
                        if col_idx < len(table_data.alignment):
                            align = table_data.alignment[col_idx]
                            if align == 'center':
                                paragraph.alignment = PP_ALIGN.CENTER
                            elif align == 'right':
                                paragraph.alignment = PP_ALIGN.RIGHT
                            else:
                                paragraph.alignment = PP_ALIGN.LEFT

                    # 数据单元格背景透明
                    fill = cell.fill
                    fill.background()  # 设置为背景色（透明）

                    # 注意：python-pptx不支持直接设置单元格边框属性
                    # 如需边框效果，建议使用表格外框或背景色区分

        # 返回表格高度（浮点数英寸值）
        return height_inches

    def _download_image(self, url: str) -> Optional[str]:
        """
        下载图片到临时目录

        Args:
            url: 图片URL

        Returns:
            本地文件路径或None
        """
        try:
            # 创建缓存目录
            cache_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'image_cache')
            os.makedirs(cache_dir, exist_ok=True)

            # 生成文件名
            filename = re.sub(r'[^\w\-_\.]', '_', url.split('/')[-1]) or 'image.png'
            local_path = os.path.join(cache_dir, filename)

            # 如果已存在,直接返回
            if os.path.exists(local_path):
                return local_path

            # 下载图片
            print(f"  下载图片: {url}")
            urllib.request.urlretrieve(url, local_path)
            return local_path

        except Exception as e:
            print(f"  警告: 图片下载失败 {url}: {e}")
            return None

    def _add_image(self, slide, img_path: str, left, top, max_width) -> float:
        """
        添加图片到幻灯片

        Args:
            slide: 幻灯片对象
            img_path: 图片路径或URL
            left: 左边距
            top: 顶部位置
            max_width: 最大宽度

        Returns:
            图片实际高度
        """
        # 检查是否为URL
        if img_path.startswith('http://') or img_path.startswith('https://'):
            local_path = self._download_image(img_path)
            if not local_path:
                return 0.0
            img_path = local_path

        # 检查文件是否存在
        if not os.path.exists(img_path):
            print(f"  警告: 图片不存在: {img_path}")
            return 0.0

        try:
            # 添加图片
            pic = slide.shapes.add_picture(img_path, left, top)

            # 计算缩放比例
            scale = min(max_width / pic.width, 1.0)
            new_width = int(pic.width * scale)
            new_height = int(pic.height * scale)

            # 调整图片大小
            pic.width = new_width
            pic.height = new_height

            # 返回英寸高度
            height_inches = new_height / 914400  # EMU to inches
            return height_inches

        except Exception as e:
            print(f"  警告: 图片插入失败: {e}")
            return 0.0

    def _add_chart(self, slide, chart_data: ChartData, left: float, top: float,
                   width: float, height: float) -> float:
        """
        添加原生图表到幻灯片

        Args:
            slide: 幻灯片对象
            chart_data: 图表数据
            left, top: 位置（英寸）
            width, height: 尺寸（英寸）

        Returns:
            实际占用高度（英寸）
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        # 映射图表类型
        type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'pie': XL_CHART_TYPE.PIE,
            'line': XL_CHART_TYPE.LINE
        }

        chart_type = type_map.get(chart_data.chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

        # 准备数据
        pptx_chart_data = CategoryChartData()
        pptx_chart_data.categories = chart_data.categories

        for series in chart_data.series:
            pptx_chart_data.add_series(series.name, series.values)

        # 添加图表
        chart_shape = slide.shapes.add_chart(
            chart_type, Inches(left), Inches(top),
            Inches(width), Inches(height), pptx_chart_data
        )
        chart = chart_shape.chart

        # 设置标题
        if chart_data.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_data.title

        # 显示图例
        chart.has_legend = True
        chart.legend.include_in_layout = False

        return height

    def add_title_slide(self, title: str, subtitle: str = ""):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)

        self.set_background(slide)

        # 计算可用区域（统一使用浮点数，单位：英寸）
        margin_left = self.config.margin_left
        margin_right = self.config.margin_right
        margin_top = self.config.margin_top
        margin_bottom = self.config.margin_bottom
        
        slide_width = self.prs.slide_width.inches if hasattr(self.prs.slide_width, 'inches') else (self.prs.slide_width / 914400)
        slide_height = self.prs.slide_height.inches if hasattr(self.prs.slide_height, 'inches') else (self.prs.slide_height / 914400)
        
        available_width = slide_width - margin_left - margin_right

        # 添加标题
        title_height = 1.5
        title_box = slide.shapes.add_textbox(
            Inches(margin_left), Inches(margin_top),
            Inches(available_width), Inches(title_height)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = title_frame.paragraphs[0]
        p.text = title
        p.font.name = self.config.title_font
        p.font.size = Pt(self.config.title_size)
        p.font.bold = True
        p.font.color.rgb = safe_color(self.config.primary_color, "#FFFFFF")
        p.alignment = PP_ALIGN.CENTER

        # 添加副标题
        if subtitle:
            subtitle_top = margin_top + title_height + 0.3
            subtitle_height = 1.0
            subtitle_box = slide.shapes.add_textbox(
                Inches(margin_left), Inches(subtitle_top),
                Inches(available_width), Inches(subtitle_height)
            )
            sub_frame = subtitle_box.text_frame
            sub_frame.word_wrap = True

            p = sub_frame.paragraphs[0]
            p.text = subtitle
            p.font.name = self.config.content_font
            p.font.size = Pt(self.config.subtitle_size)
            p.font.color.rgb = safe_color(self.config.secondary_color, "#E0E0E0")
            p.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, slide_content: SlideContent):
        """添加内容页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)

        self.set_background(slide)

        # 计算可用区域（统一使用浮点数，单位：英寸）
        margin_left = float(self.config.margin_left)
        margin_right = float(self.config.margin_right)
        margin_top = float(self.config.margin_top)
        margin_bottom = float(self.config.margin_bottom)

        # 获取幻灯片尺寸（转换为浮点数英寸）
        slide_width = float(self.prs.slide_width.inches) if hasattr(self.prs.slide_width, 'inches') else float(self.prs.slide_width / 914400)
        slide_height = float(self.prs.slide_height.inches) if hasattr(self.prs.slide_height, 'inches') else float(self.prs.slide_height / 914400)

        available_width = slide_width - margin_left - margin_right
        available_height = slide_height - margin_top - margin_bottom

        # 安全检查：确保可用高度为正数
        if available_height <= 0:
            print(f"警告：页面空间不足 (available_height={available_height:.2f}\")")
            return

        current_top = margin_top  # 浮点数

        # 添加标题
        title_height = 0.8  # 英寸（浮点数）
        title_box = slide.shapes.add_textbox(
            Inches(margin_left), Inches(current_top),
            Inches(available_width), Inches(title_height)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        p = title_frame.paragraphs[0]
        p.text = slide_content.title
        p.font.name = self.config.title_font
        p.font.size = Pt(self.config.title_size)
        p.font.bold = True
        p.font.color.rgb = safe_color(self.config.primary_color, "#FFFFFF")
        p.alignment = PP_ALIGN.LEFT

        current_top += title_height + 0.2

        # 添加副标题
        if slide_content.subtitle:
            subtitle_height = 0.5
            subtitle_box = slide.shapes.add_textbox(
                Inches(margin_left), Inches(current_top),
                Inches(available_width), Inches(subtitle_height)
            )
            sub_frame = subtitle_box.text_frame
            sub_frame.word_wrap = True

            p = sub_frame.paragraphs[0]
            p.text = slide_content.subtitle
            p.font.name = self.config.content_font
            p.font.size = Pt(self.config.subtitle_size - 4)
            p.font.color.rgb = safe_color(self.config.secondary_color, "#E0E0E0")
            p.alignment = PP_ALIGN.LEFT

            current_top += subtitle_height + 0.2

        # 添加图片
        if slide_content.images:
            for img_path in slide_content.images:
                img_height = self._add_image(slide, img_path, margin_left, current_top, available_width)
                current_top += img_height + 0.3

        # 添加表格
        if slide_content.table:
            table_top = current_top
            table_height = self._add_table(slide, slide_content.table, margin_left, table_top, available_width)
            # 确保table_height是数值类型
            if not isinstance(table_height, (int, float)):
                print(f"警告: _add_table返回了非数值类型: {type(table_height)}")
                table_height = 0.0
            current_top += table_height + 0.3

        # 添加项目符号列表
        if slide_content.bullet_points:
            # 计算剩余可用高度（浮点数）
            remaining_height = available_height - (current_top - margin_top)

            # 如果剩余空间不足，跳过
            if remaining_height < 1.0:
                print(f"警告: 页面空间不足，跳过 {len(slide_content.bullet_points)} 个列表项")
            else:
                bullet_box = slide.shapes.add_textbox(
                    Inches(margin_left), Inches(current_top),
                    Inches(available_width), Inches(remaining_height)
                )
                bullet_frame = bullet_box.text_frame
                bullet_frame.word_wrap = True
                bullet_frame.line_spacing = self.config.line_spacing

                for i, point in enumerate(slide_content.bullet_points):
                    if i == 0:
                        p = bullet_frame.paragraphs[0]
                    else:
                        p = bullet_frame.add_paragraph()

                    # 检测是否为二级目录项（以 • 开头且前面有空格）
                    is_subitem = point.strip().startswith('•') or point.startswith('   ')

                    # 应用Markdown加粗格式
                    apply_bold_formatting(p, point.strip())

                    p.font.name = self.config.content_font

                    # 二级标题使用较小字体
                    if is_subitem:
                        p.font.size = Pt(self.config.body_size - 4)
                        p.level = 1  # 缩进级别
                    else:
                        p.font.size = Pt(self.config.body_size)
                        p.level = 0

                    p.font.color.rgb = safe_color(self.config.text_color, "#FFFFFF")
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)

                # 更新当前位置（估算列表占用高度，浮点数）
                list_height = len(slide_content.bullet_points) * (self.config.body_size / 72.0 * self.config.line_spacing + 0.1)
                current_top += list_height

        # 添加普通文本内容
        if slide_content.content_lines:
            # 合并内容行
            full_text = ' '.join(slide_content.content_lines)

            # 计算文本实际占用高度
            text_height = calculate_text_height(
                full_text,
                self.config.body_size,
                self.config.line_spacing,
                available_width  # 已经是浮点数
            )

            # 计算剩余可用高度（浮点数）
            remaining_height = available_height - (current_top - margin_top)

            # 如果文本会超出页面，截断并给出警告
            if text_height > remaining_height:
                print(f"警告: 文本内容超出页面空间 ({text_height:.2f}\" > {remaining_height:.2f}\")")
                # 仍然添加，但限制高度
                text_height_to_use = min(text_height, remaining_height)
            else:
                text_height_to_use = text_height

            content_box = slide.shapes.add_textbox(
                Inches(margin_left), Inches(current_top),
                Inches(available_width), Inches(text_height_to_use)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.line_spacing = self.config.line_spacing

            # 智能分段
            paragraphs = self._split_text_into_paragraphs(full_text)

            for i, paragraph in enumerate(paragraphs):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()

                # 应用Markdown加粗格式
                apply_bold_formatting(p, paragraph)

                # 设置字体属性（注意：apply_bold_formatting可能已经设置了部分run的bold）
                p.font.name = self.config.content_font
                p.font.size = Pt(self.config.body_size)
                p.font.color.rgb = safe_color(self.config.text_color, "#FFFFFF")
                p.space_before = Pt(6)
                p.space_after = Pt(6)

            # 更新当前位置（浮点数）
            current_top += text_height_to_use

        # 添加代码块
        if slide_content.code_blocks:
            for code_block in slide_content.code_blocks:
                remaining_height = available_height - (current_top - margin_top)

                if remaining_height < 1.0:
                    print(f"警告: 页面空间不足，跳过代码块")
                    break

                # 检测是否为Mermaid图表
                if code_block.language.lower() in ['mermaid', 'mmd'] or code_block.is_diagram:
                    # 尝试转换为图片
                    from .mermaid_converter import convert_mermaid_to_image

                    img_path = convert_mermaid_to_image(code_block.content, output_format='png')

                    if img_path:
                        # 作为图片插入
                        img_height = self._add_image(slide, img_path, margin_left, current_top, available_width)
                        current_top += img_height + 0.2
                        continue
                    else:
                        # 降级为代码块显示
                        print("  降级: Mermaid图表将以代码形式显示")

                # 普通代码块渲染（原有逻辑）
                code_lines_count = len(code_block.content.split('\n'))
                code_height = code_lines_count * (12 / 72.0 * 1.1) + 0.3

                if code_height > remaining_height:
                    print(f"警告: 代码块超出页面空间")
                    code_height = remaining_height

                # 创建代码框
                code_box = slide.shapes.add_textbox(
                    Inches(margin_left), Inches(current_top),
                    Inches(available_width), Inches(code_height)
                )
                code_frame = code_box.text_frame
                code_frame.word_wrap = True
                code_frame.line_spacing = 1.1

                # 设置背景色（深色主题）
                code_box.fill.solid()
                code_box.fill.fore_color.rgb = safe_color("#2D2D2D", "#F5F5F5")

                # 添加代码内容
                code_text = code_block.content
                p = code_frame.paragraphs[0]
                p.text = code_text

                # 设置等宽字体
                for paragraph in code_frame.paragraphs:
                    paragraph.font.name = "Consolas" if os.name == 'nt' else "Monaco"
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = safe_color("#E0E0E0", "#FFFFFF")

                current_top += code_height + 0.2

        # 添加图表
        if slide_content.charts:
            for chart_data in slide_content.charts:
                remaining_height = available_height - (current_top - margin_top)

                if remaining_height < 3.0:  # 图表最小高度
                    print(f"警告: 页面空间不足，跳过图表")
                    break

                chart_height = min(self.config.default_chart_height, remaining_height)
                chart_width = self.config.default_chart_width

                self._add_chart(slide, chart_data, margin_left, current_top,
                               chart_width, chart_height)
                current_top += chart_height + 0.3

    def _split_text_into_paragraphs(self, text: str) -> List[str]:
        """
        将长文本智能分割为多个段落
        考虑每行最大字符数和每页最大行数
        """
        max_chars = self.config.max_chars_per_line
        max_lines = self.config.max_lines_per_slide

        # 按句子分割 - 修复Python 3.6兼容性,避免空模式匹配
        # 使用非空断言确保正则表达式有效
        sentences = re.split(r'(?<=[。！？.!?])\s+', text) if text else []

        paragraphs = []
        current_paragraph = ""
        line_count = 0

        for sentence in sentences:
            if not sentence.strip():
                continue

            # 计算这句话需要的行数
            lines_needed = (len(sentence) + max_chars - 1) // max_chars

            # 如果当前段落加上新句子会超过限制,开始新段落
            if line_count + lines_needed > max_lines and current_paragraph:
                paragraphs.append(current_paragraph)
                current_paragraph = sentence
                line_count = lines_needed
            else:
                if current_paragraph:
                    current_paragraph += " " + sentence
                else:
                    current_paragraph = sentence
                line_count += lines_needed

        # 添加最后一个段落
        if current_paragraph:
            paragraphs.append(current_paragraph)

        return paragraphs if paragraphs else [text]

    def generate(self, content: PresentationContent, output_path: str):
        """
        生成PPT文件

        Args:
            content: 结构化的演示文稿内容
            output_path: 输出文件路径
        """
        # 清空现有幻灯片
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

        total_slides = len(content.slides)

        # 添加标题页
        if content.slides:
            first_slide = content.slides[0]
            self.add_title_slide(first_slide.title, first_slide.subtitle)
            # 添加页脚到标题页(如果启用)
            if self.config.show_footer_on_title_slide:
                self._add_footer(self.prs.slides[0], 1, total_slides)

            # 添加内容页
            for idx, slide_content in enumerate(content.slides[1:], 2):
                self.add_content_slide(slide_content)
                # 添加页脚
                self._add_footer(self.prs.slides[idx - 1], idx, total_slides)

        # 保存文件
        self.prs.save(output_path)
        print(f"PPT已生成: {output_path}")
        print(f"共 {len(self.prs.slides)} 页")

    def generate_with_smart_layout(self, content: PresentationContent, output_path: str):
        """
        使用智能布局生成PPT（两阶段方法）

        阶段1: 内容预扫描 - 计算每个内容块的实际高度
        阶段2: 智能分页 - 根据内容高度分配到页面
        阶段3: 按分页结果渲染 - 避免重叠

        Args:
            content: 结构化的演示文稿内容
            output_path: 输出文件路径
        """
        # 清空现有幻灯片
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

        print("\n" + "="*60)
        print("智能布局模式 - 开始内容预扫描")
        print("="*60)

        # 阶段 1 & 2: 创建布局计划（预扫描 + 智能分页）
        layout_plan = create_layout_plan(content.slides, self.config)
        total_pages = len(layout_plan)

        print(f"\n预扫描完成: 共 {total_pages} 页内容")

        # 添加标题页（使用第一页的内容）
        if content.slides:
            first_slide = content.slides[0]
            self.add_title_slide(first_slide.title, first_slide.subtitle)
            if self.config.show_footer_on_title_slide:
                self._add_footer(self.prs.slides[0], 1, total_pages + 1)

        # 阶段 3: 渲染每个内容页
        for page_idx, page_blocks in enumerate(layout_plan, 1):
            print(f"\n渲染第 {page_idx}/{total_pages} 页...")
            self._render_content_page(page_blocks)
            
            # 添加页脚
            self._add_footer(self.prs.slides[page_idx], page_idx + 1, total_pages + 1)

        # 保存文件
        self.prs.save(output_path)
        print(f"\n{'='*60}")
        print(f"PPT已生成: {output_path}")
        print(f"共 {len(self.prs.slides)} 页")
        print(f"{'='*60}\n")

    def _render_content_page(self, blocks: List[ContentBlock]):
        """
        渲染单个内容页（使用预计算的内容块）

        Args:
            blocks: 该页的内容块列表
        """
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        self.set_background(slide)

        # 计算可用区域
        margin_left = float(self.config.margin_left)
        margin_right = float(self.config.margin_right)
        margin_top = float(self.config.margin_top)
        margin_bottom = float(self.config.margin_bottom)

        slide_width = float(self.prs.slide_width.inches) if hasattr(self.prs.slide_width, 'inches') else float(self.prs.slide_width / 914400)
        slide_height = float(self.prs.slide_height.inches) if hasattr(self.prs.slide_height, 'inches') else float(self.prs.slide_height / 914400)

        available_width = slide_width - margin_left - margin_right
        available_height = slide_height - margin_top - margin_bottom

        current_top = margin_top

        # 按顺序渲染每个内容块
        for block in blocks:
            # 检查空间是否足够
            remaining_height = available_height - (current_top - margin_top)
            
            if remaining_height < block.min_height and block.min_height > 0:
                print(f"  警告: 空间不足，跳过 {block.type} 块 (需要 {block.height:.2f}\", 剩余 {remaining_height:.2f}\")")
                continue

            if block.type == 'title':
                # 渲染标题
                title_box = slide.shapes.add_textbox(
                    Inches(margin_left), Inches(current_top),
                    Inches(available_width), Inches(block.height)
                )
                title_frame = title_box.text_frame
                title_frame.word_wrap = True

                p = title_frame.paragraphs[0]
                p.text = block.content
                p.font.name = self.config.title_font
                p.font.size = Pt(self.config.title_size)
                p.font.bold = True
                p.font.color.rgb = safe_color(self.config.primary_color, "#FFFFFF")
                p.alignment = PP_ALIGN.LEFT

            elif block.type == 'subtitle':
                # 渲染副标题
                subtitle_box = slide.shapes.add_textbox(
                    Inches(margin_left), Inches(current_top),
                    Inches(available_width), Inches(block.height)
                )
                sub_frame = subtitle_box.text_frame
                sub_frame.word_wrap = True

                p = sub_frame.paragraphs[0]
                p.text = block.content
                p.font.name = self.config.content_font
                p.font.size = Pt(self.config.subtitle_size - 4)
                p.font.color.rgb = safe_color(self.config.secondary_color, "#E0E0E0")
                p.alignment = PP_ALIGN.LEFT

            elif block.type == 'image':
                # 渲染图片
                img_path = block.content['path']
                img_height = self._add_image(slide, img_path, margin_left, current_top, available_width)
                # 更新实际高度（可能与预估不同）
                block.height = img_height + 0.3

            elif block.type == 'table':
                # 渲染表格
                table_data = block.content
                table_height = self._add_table(slide, table_data, margin_left, current_top, available_width)
                block.height = table_height + 0.3

            elif block.type == 'bullet':
                # 渲染项目符号列表
                bullet_points = block.content
                bullet_box = slide.shapes.add_textbox(
                    Inches(margin_left), Inches(current_top),
                    Inches(available_width), Inches(block.height)
                )
                bullet_frame = bullet_box.text_frame
                bullet_frame.word_wrap = True
                bullet_frame.line_spacing = self.config.line_spacing

                for i, point in enumerate(bullet_points):
                    if i == 0:
                        p = bullet_frame.paragraphs[0]
                    else:
                        p = bullet_frame.add_paragraph()

                    is_subitem = point.strip().startswith('•') or point.startswith('   ')
                    apply_bold_formatting(p, point.strip())

                    p.font.name = self.config.content_font
                    if is_subitem:
                        p.font.size = Pt(self.config.body_size - 4)
                        p.level = 1
                    else:
                        p.font.size = Pt(self.config.body_size)
                        p.level = 0

                    p.font.color.rgb = safe_color(self.config.text_color, "#FFFFFF")
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)

            elif block.type == 'text':
                # 渲染普通文本
                full_text = block.content
                text_box = slide.shapes.add_textbox(
                    Inches(margin_left), Inches(current_top),
                    Inches(available_width), Inches(block.height)
                )
                content_frame = text_box.text_frame
                content_frame.word_wrap = True
                content_frame.line_spacing = self.config.line_spacing

                paragraphs = self._split_text_into_paragraphs(full_text)
                for i, paragraph in enumerate(paragraphs):
                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()

                    apply_bold_formatting(p, paragraph)
                    p.font.name = self.config.content_font
                    p.font.size = Pt(self.config.body_size)
                    p.font.color.rgb = safe_color(self.config.text_color, "#FFFFFF")
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)

            elif block.type == 'code':
                # 渲染代码块
                code_block = block.content
                code_box = slide.shapes.add_textbox(
                    Inches(margin_left), Inches(current_top),
                    Inches(available_width), Inches(block.height)
                )
                code_frame = code_box.text_frame
                code_frame.word_wrap = True
                code_frame.line_spacing = 1.1

                code_box.fill.solid()
                code_box.fill.fore_color.rgb = safe_color("#2D2D2D", "#F5F5F5")

                p = code_frame.paragraphs[0]
                p.text = code_block.content

                for paragraph in code_frame.paragraphs:
                    paragraph.font.name = "Consolas" if os.name == 'nt' else "Monaco"
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = safe_color("#E0E0E0", "#FFFFFF")

            elif block.type == 'mermaid':
                # 渲染 Mermaid 图表
                from .mermaid_converter import convert_mermaid_to_image

                code_content = block.content['code']
                img_path = convert_mermaid_to_image(code_content, output_format='png')

                if img_path:
                    # 作为图片插入
                    img_height = self._add_image(slide, img_path, margin_left, current_top, available_width)
                    block.height = img_height + 0.2
                else:
                    # 降级为代码块显示
                    print("  降级: Mermaid图表将以代码形式显示")
                    code_box = slide.shapes.add_textbox(
                        Inches(margin_left), Inches(current_top),
                        Inches(available_width), Inches(block.height)
                    )
                    code_frame = code_box.text_frame
                    code_frame.word_wrap = True

                    code_box.fill.solid()
                    code_box.fill.fore_color.rgb = safe_color("#2D2D2D", "#F5F5F5")

                    p = code_frame.paragraphs[0]
                    p.text = code_content

                    for paragraph in code_frame.paragraphs:
                        paragraph.font.name = "Consolas" if os.name == 'nt' else "Monaco"
                        paragraph.font.size = Pt(12)
                        paragraph.font.color.rgb = safe_color("#E0E0E0", "#FFFFFF")

            elif block.type == 'chart':
                # 渲染原生图表
                chart_data = block.content
                chart_height = self._add_chart(slide, chart_data, margin_left, current_top,
                                              available_width, block.height - 0.3)
                block.height = chart_height + 0.3

            # 更新当前位置（使用块的实际或预估高度）
            current_top += block.height

            print(f"  ✓ 渲染 {block.type} (高度: {block.height:.2f}\")")

    def _add_footer(self, slide, current_page: int, total_pages: int):
        """
        添加页脚和页码到幻灯片

        Args:
            slide: 幻灯片对象
            current_page: 当前页码
            total_pages: 总页数
        """
        if not self.config.show_page_number and not self.config.footer_text:
            return

        # 构建页脚文本
        footer_parts = []

        # 添加自定义页脚文本
        if self.config.footer_text:
            footer_parts.append(self.config.footer_text)

        # 添加页码
        if self.config.show_page_number:
            page_text = self.config.page_number_format.format(
                current=current_page,
                total=total_pages
            )
            footer_parts.append(page_text)

        if not footer_parts:
            return

        footer_text = " | ".join(footer_parts)

        # 计算页脚位置（使用浮点数）
        margin_bottom = self.config.margin_bottom
        footer_height = 0.4
        
        slide_height = self.prs.slide_height.inches if hasattr(self.prs.slide_height, 'inches') else (self.prs.slide_height / 914400)
        footer_top = slide_height - margin_bottom - footer_height

        # 创建页脚文本框
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(footer_top),
            Inches(slide_height - 1.0), Inches(footer_height)
        )
        footer_frame = footer_box.text_frame
        footer_frame.word_wrap = True

        p = footer_frame.paragraphs[0]
        p.text = footer_text
        p.font.name = self.config.content_font
        p.font.size = Pt(self.config.footer_font_size)
        p.font.color.rgb = safe_color(self.config.footer_color, "#666666")
        p.alignment = PP_ALIGN.CENTER


# 需要导入re模块
import re
